from typing import List, Dict, Optional, Tuple
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw
import zipfile
import random

from llm_router import ProviderConfig, make_outline_with_llm

MAX_SLIDE_PREVIEWS = 12

def _collect_template_images(template_bytes: bytes) -> List[bytes]:
    imgs = []
    with zipfile.ZipFile(BytesIO(template_bytes), 'r') as zf:
        for name in zf.namelist():
            if name.startswith("ppt/media/") and name.lower().split('.')[-1] in {"png","jpg","jpeg"}:
                try:
                    imgs.append(zf.read(name))
                except Exception:
                    continue
    return imgs

def _add_image_if_layout_allows(slide, img_bytes: bytes):
    try:
        pic_stream = BytesIO(img_bytes)
        # place on right side with a max width heuristic
        slide.shapes.add_picture(pic_stream, Inches(6.0), Inches(1.5), width=Inches(3.0))
    except Exception:
        pass

def _apply_bullets(text_frame, bullets: List[str]):
    # Clear existing text safely
    try:
        for p in list(text_frame.paragraphs):
            p.clear()
    except Exception:
        pass
    if not bullets:
        return
    for i, b in enumerate(bullets):
        p = text_frame.add_paragraph() if i>0 else text_frame.paragraphs[0]
        p.text = b
        p.level = 0
        p.space_after = Pt(4)

def _title_content_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        names = (layout.name or "").lower()
        if "title" in names and ("content" in names or "text" in names or "body" in names):
            return layout
    return prs.slide_layouts[0]

def _title_only_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        if "title only" in (layout.name or "").lower():
            return layout
    return prs.slide_layouts[0]

def _make_cover_slide(prs: Presentation, title: str, subtitle: Optional[str] = None):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    try:
        if slide.shapes.title:
            slide.shapes.title.text = title
    except Exception:
        pass
    try:
        for ph in slide.placeholders:
            if getattr(ph, "placeholder_format", None) and getattr(ph.placeholder_format, "idx", None) == 1:
                ph.text = subtitle or ""
    except Exception:
        pass
    return slide

def build_presentation(text: str, outline: List[Dict], template_stream, guidance: Optional[str]=None,
                       provider_cfg: Optional[ProviderConfig]=None, preview: bool=False) -> Tuple[bytes, Optional[List]]:
    template_bytes = template_stream.read()
    prs = Presentation(BytesIO(template_bytes))

    template_imgs = _collect_template_images(template_bytes)
    random.shuffle(template_imgs)

    first_title = outline[0]["title"] if outline else "Presentation"
    subtitle = (guidance or "").strip().capitalize() if guidance else None
    _make_cover_slide(prs, first_title, subtitle)

    layout_tc = _title_content_layout(prs)
    layout_title_only = _title_only_layout(prs)

    img_idx = 0
    for i, slide_spec in enumerate(outline, start=1):
        slide = prs.slides.add_slide(layout_tc)
        try:
            if slide.shapes.title:
                slide.shapes.title.text = slide_spec.get("title","Slide")
        except Exception:
            pass

        # find body placeholder
        body_placeholder = None
        for shape in slide.placeholders:
            try:
                if shape.has_text_frame and getattr(shape, 'placeholder_format', None) and shape.placeholder_format.idx != 0:
                    body_placeholder = shape
                    break
            except Exception:
                continue
        if body_placeholder is None:
            slide = prs.slides.add_slide(layout_title_only)
            try:
                if slide.shapes.title:
                    slide.shapes.title.text = slide_spec.get("title","Slide")
            except Exception:
                pass
        else:
            _apply_bullets(body_placeholder.text_frame, slide_spec.get("bullets", []))

        if template_imgs:
            _add_image_if_layout_allows(slide, template_imgs[img_idx % len(template_imgs)])
            img_idx += 1

        if provider_cfg is not None:
            try:
                notes = slide.notes_slide.notes_text_frame
                notes.text = _gen_notes(slide_spec, provider_cfg)
            except Exception:
                pass

    bio = BytesIO()
    prs.save(bio)
    pptx_bytes = bio.getvalue()

    previews = []
    if preview:
        # Text-based thumbnails (safe, quick)
        for spec in outline[:MAX_SLIDE_PREVIEWS]:
            img = Image.new("RGB", (1024, 576), (245, 245, 245))
            d = ImageDraw.Draw(img)
            title = spec.get("title","Slide")
            # wrap title
            lines = []
            w = 34
            for i in range(0, len(title), w):
                lines.append(title[i:i+w])
            y = 24
            d.text((30, y), "\n".join(lines), fill=(30,30,30))
            # bullets
            y += 100
            for j, b in enumerate(spec.get("bullets", [])[:6]):
                tb = b if len(b) <= 80 else b[:77] + '...'
                d.text((40, y + j*28), f"• {tb}", fill=(60,60,60))
            previews.append(img)

    return pptx_bytes, previews

def _gen_notes(slide_spec: Dict, cfg: ProviderConfig) -> str:
    prompt = f"Write concise speaker notes for a slide titled \"{slide_spec.get('title','')}\".\nBullets:\n- " + "\n- ".join(slide_spec.get("bullets", [])[:6])
    try:
        outline = make_outline_with_llm(prompt, "Write a short paragraph of 60-100 words.", cfg, max_slides=1)
        if isinstance(outline, list) and outline:
            bs = outline[0].get("bullets", [])
            if bs:
                return " ".join(bs)
    except Exception:
        pass
    # Fallback: join bullets
    return " ".join(slide_spec.get("bullets", []))
